"""The coach's/business owner's uploaded knowledge base - documents (PDFs,
photos, presentations, via Docling) ingested and made searchable, global,
not scoped to any one contact. Local Ollama embeddings (nomic-embed-text)
backed by the Qdrant instance Adiyan already runs against.

Per-contact conversation memory used to live in this same file, backed by a
plain LlamaIndex VectorStoreIndex - it moved to mesh/memory/mem0_backend.py,
which wraps mem0ai instead: a real extraction+consolidation pipeline
(atomic facts, ADD/merge against existing memories), not a second hand-rolled
memory engine duplicating what this file already does for documents."""
import base64
import io
import json
import logging
import re
import sqlite3
from pathlib import Path
from typing import Any, Dict, List, Optional

from llama_index.core import Document, VectorStoreIndex
from llama_index.core.node_parser import MarkdownNodeParser, SentenceSplitter
from llama_index.core.vector_stores import FilterOperator, MetadataFilter, MetadataFilters
from llama_index.embeddings.ollama import OllamaEmbedding
from llama_index.vector_stores.qdrant import QdrantVectorStore
from qdrant_client import QdrantClient
from qdrant_client.models import FieldCondition, Filter, MatchValue

from mesh.lib.paths import kb_documents_dir, state_db_path
from mesh.memory.constants import AGENT_ID

logger = logging.getLogger('MemoryIndex')

DOCUMENTS_TABLE_SCHEMA = """
CREATE TABLE IF NOT EXISTS kb_documents (
    source_filename TEXT PRIMARY KEY,
    stored_path TEXT NOT NULL,
    mimetype TEXT,
    ingested_at TEXT NOT NULL,
    chunks INTEGER NOT NULL
)
"""


def _safe_filename(filename: str) -> str:
    """Strips any directory component and anything that isn't safe as a
    bare filename - filename arrives from a WhatsApp upload's own name/
    caption field, not something to trust blindly as a path component."""
    name = Path(filename).name
    return re.sub(r'[^A-Za-z0-9._-]', '_', name) or 'document'


def _documents_db() -> sqlite3.Connection:
    conn = sqlite3.connect(state_db_path(AGENT_ID))
    conn.execute(DOCUMENTS_TABLE_SCHEMA)
    return conn


def _extract_pptx_markdown(content: bytes, filename: str) -> str:
    """Docling's own PPTX backend only extracts native text shapes and
    represents embedded pictures as bare '<!-- image -->' placeholders - it
    does not run OCR on them the way its image/PDF pipeline does. Confirmed
    live: a 13-slide deck where every slide was a single full-slide picture
    (a common export shape from slide-design tools) round-tripped through
    the normal ingest path as twelve literal "<!-- image -->" lines and
    nothing else - a spelling-mistake analysis had no real text to work
    with at all, and reasonably reported finding none.

    Works around this by handling PPTX specially: for each slide, prefer
    its own real text shapes/tables if any exist; when a slide has none
    (the all-picture case), OCR its picture shapes individually through
    Docling's image pipeline instead - the same path that already works
    correctly for a standalone photo/screenshot (this same ingest_document's
    non-PPTX branch, proven live against an Aadhaar card)."""
    from docling.document_converter import DocumentConverter
    from docling_core.types.io import DocumentStream
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def _collect_text(shapes) -> List[str]:
        texts = []
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                texts.extend(_collect_text(shape.shapes))
            elif shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            texts.append(cell.text)
        return texts

    def _collect_pictures(shapes) -> List:
        pictures = []
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                pictures.extend(_collect_pictures(shape.shapes))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures.append(shape.image)
        return pictures

    prs = Presentation(io.BytesIO(content))
    converter = DocumentConverter()
    slide_blocks = []

    for i, slide in enumerate(prs.slides):
        texts = _collect_text(slide.shapes)
        if texts:
            slide_blocks.append(f'## Slide {i + 1}\n\n' + '\n\n'.join(texts))
            continue

        ocr_texts = []
        for image in _collect_pictures(slide.shapes):
            try:
                result = converter.convert(
                    DocumentStream(name=f'slide_{i + 1}.{image.ext}', stream=io.BytesIO(image.blob))
                )
                ocr_text = result.document.export_to_markdown().strip()
                if ocr_text and ocr_text != '<!-- image -->':
                    ocr_texts.append(ocr_text)
            except Exception as e:
                logger.warning(f"Could not OCR slide {i + 1} of {filename!r}: {e}")
        if ocr_texts:
            slide_blocks.append(f'## Slide {i + 1}\n\n' + '\n\n'.join(ocr_texts))
        else:
            # Not an exception - Docling's OCR ran and simply found nothing
            # readable in this slide's picture(s) (confirmed live: happens
            # even for a slide with a perfectly legible title, an inherent
            # OCR reliability limit, not a bug in this loop). Logged, not
            # silently dropped - a coach reviewing why their deck's
            # analysis missed a slide should be able to find out why.
            logger.warning(f"No extractable text from slide {i + 1} of {filename!r} - dropped from the ingested text")

    return '\n\n'.join(slide_blocks)

# New name, not the old 'coaching_history' the placeholder-vector code used to write to
# (that collection held only fake constant vectors - nothing worth migrating from it).
KB_COLLECTION_NAME = 'adiyan_knowledge_base'
KB_DEFAULT_TOP_K = 4
KB_CHUNK_SIZE = 800
KB_CHUNK_OVERLAP = 100

# Separate collection from KB_COLLECTION_NAME on purpose - a page is
# addressed by exact page_number (Scheduler's "send tonight's page"), never
# semantically searched the way knowledge-base chunks are, so it has no
# business sharing a collection (or a chunking strategy) with content that
# genuinely is searched. Each point still carries a real embedding (nothing
# about VectorStoreIndex.insert() supports skipping that), it's just never
# queried by similarity.
KB_PAGES_COLLECTION_NAME = 'adiyan_book_pages'

# search_within_document()'s own default - a bit higher than KB_DEFAULT_TOP_K
# since the search space here is already narrowed to one document, so
# returning a few more candidate chunks costs little and raises the odds the
# actually-relevant one is among them.
DOC_SEARCH_DEFAULT_TOP_K = 5

# find_source_document()'s cutoff below which a top-1 "match" is treated as
# no match at all, rather than confidently returning the least-bad option in
# a knowledge base that has nothing genuinely relevant. First-pass value, not
# empirically tuned across many real queries yet - nomic-embed-text cosine
# similarity for a genuinely relevant chunk has run well above this in
# testing so far, and the specific failure this exists for (a Vizag-trip
# query "matching" an Atomic Habits PDF) scored well below it. Revisit if
# this starts rejecting real matches or letting weak ones through.
SOURCE_MATCH_MIN_SCORE = 0.55

_instances: Dict[tuple, 'MemoryIndex'] = {}


class MemoryIndex:
    """The knowledge-base side only - a LlamaIndex VectorStoreIndex over Qdrant
    for coach/owner-uploaded documents. Conversation memory lives in
    mesh/memory/mem0_backend.py now, not here - see this module's own
    docstring."""

    def __init__(self, qdrant_url: str, ollama_url: str, embed_model_name: str = 'nomic-embed-text'):
        self.embed_model = OllamaEmbedding(model_name=embed_model_name, base_url=ollama_url)
        client = QdrantClient(url=qdrant_url)
        self._qdrant_client = client

        kb_vector_store = QdrantVectorStore(client=client, collection_name=KB_COLLECTION_NAME)
        self.kb_index = VectorStoreIndex.from_vector_store(kb_vector_store, embed_model=self.embed_model)
        self._splitter = SentenceSplitter(chunk_size=KB_CHUNK_SIZE, chunk_overlap=KB_CHUNK_OVERLAP)
        self._markdown_splitter = MarkdownNodeParser()

        pages_vector_store = QdrantVectorStore(client=client, collection_name=KB_PAGES_COLLECTION_NAME)
        self.pages_index = VectorStoreIndex.from_vector_store(pages_vector_store, embed_model=self.embed_model)

    # Conversation memory (insert/retrieve) used to live here, backed by a
    # plain LlamaIndex VectorStoreIndex with zero callers of insert() ever
    # found in this repo - see mesh/memory/mem0_backend.py, which replaces
    # it with a real extraction+consolidation pipeline (mem0ai) instead of
    # this file continuing to grow a second, hand-rolled memory engine
    # alongside the knowledge-base one below.

    def _split_text(self, markdown: str, safe_name: str) -> List[str]:
        """Confirmed live: EXTERNAL_DEPENDENCIES.md's real "what port does
        Qdrant run on" eval case scored only 0.496 similarity (below
        SOURCE_MATCH_MIN_SCORE) despite the fact being right there in the
        text - traced to the heading-blind SentenceSplitter concatenating
        five unrelated ## sections (Ollama, Qdrant, MongoDB, Phoenix, the
        doc intro) into one 3000-char chunk, diluting the embedding for any
        single one of them. Isolating just the Qdrant section and
        re-embedding it alone measured 0.665 for the same query - this is
        the fix, not a threshold tweak.

        Native markdown source documents (.md/.markdown) get MarkdownNodeParser
        instead - splits on real ## heading boundaries, so each topic gets
        its own chunk. Docling-converted content (PDF/PPTX/OCR'd images)
        keeps the plain SentenceSplitter - its markdown headings come from
        OCR/extraction heuristics, not an author's real document structure,
        and are a less reliable place to split on.

        MarkdownNodeParser has no max chunk size of its own - a section
        under one heading that runs unusually long still needs
        KB_CHUNK_SIZE's own cap applied, or one oversized node could
        reintroduce the exact dilution problem this exists to fix. Each
        node is re-run through split_text() unconditionally rather than
        checked first - a node already under chunk_size comes back as a
        single unchanged chunk, so this is a safety net, not redundant
        work."""
        if not safe_name.lower().endswith(('.md', '.markdown')):
            return self._splitter.split_text(markdown)

        nodes = self._markdown_splitter.get_nodes_from_documents([Document(text=markdown)])
        chunks: List[str] = []
        for node in nodes:
            chunks.extend(self._splitter.split_text(node.get_content()))
        return chunks

    def ingest_document(
        self, content: bytes, filename: str, timestamp: str, username: str, mimetype: Optional[str] = None,
    ) -> tuple:
        """Returns (chunks_count, source_filename) - source_filename is the composite
        <username>/<filename> key this document is now stored under everywhere
        (Qdrant chunk metadata, the raw-file index), needed by a caller that wants to
        act on this exact document right after ingesting it (see
        mesh/analysis/skills/analyze.py's source_filename param, used by the
        upload+instruct combined flow in mesh/orchestrator/skills/handle_message.py to
        skip a separate resolve_document lookup). Raises on parse failure - the caller
        (mesh/memory/skills/ingest.py) is expected to report that back to the coach
        over WhatsApp, not swallow it, since this is a coach-initiated action they need
        feedback on.

        Also saves the original bytes to disk under kb_documents_dir/<username>/ and
        records them in a small local index (kb_documents table) - the knowledge base's
        vectors alone can only answer "what does it say" (retrieve_knowledge_base),
        never "send me that document back" (get_document/find_source_document), so a
        later share needs the original file kept somewhere too.

        username folds into source_filename's own identity (<username>/<filename>), not
        just the physical path - so two different uploaders using the same original
        filename never collide, in Qdrant chunk metadata and the raw-file index alike,
        not only on disk."""
        safe_user = _safe_filename(username)
        safe_name = _safe_filename(filename)
        source_key = f'{safe_user}/{safe_name}'

        if safe_name.lower().endswith(('.pptx', '.ppt')):
            # See _extract_pptx_markdown()'s own docstring - Docling's PPTX
            # backend doesn't OCR embedded pictures, which every slide in a
            # slide-design-tool export (Canva, Figma, NotebookLM, etc.)
            # typically is.
            markdown = _extract_pptx_markdown(content, safe_name)
        else:
            from docling.document_converter import DocumentConverter
            from docling_core.types.io import DocumentStream

            converter = DocumentConverter()
            result = converter.convert(DocumentStream(name=safe_name, stream=io.BytesIO(content)))
            markdown = result.document.export_to_markdown()

        if not markdown.strip():
            raise ValueError(f"Docling extracted no text from '{filename}' (scanned image with no OCR text?)")

        # Delete any chunks already on file for this exact source_key before
        # inserting fresh ones - .insert() below only adds, it doesn't
        # replace, so re-ingesting the same document (a coach re-uploading
        # it, or this same fix re-processing a document ingested before a
        # quality improvement like _extract_pptx_markdown existed) would
        # otherwise leave stale old chunks sitting alongside the new ones
        # forever, with colliding chunk_index values confusing
        # get_document_text()'s ordering.
        self._qdrant_client.delete(
            collection_name=KB_COLLECTION_NAME,
            points_selector=Filter(must=[FieldCondition(key='source_filename', match=MatchValue(value=source_key))]),
        )

        chunks = self._split_text(markdown, safe_name)
        for i, chunk in enumerate(chunks):
            self.kb_index.insert(Document(
                text=chunk,
                metadata={
                    'source_filename': source_key,
                    'chunk_index': i,
                    'ingested_at': timestamp,
                },
            ))

        user_dir = kb_documents_dir(AGENT_ID) / safe_user
        user_dir.mkdir(parents=True, exist_ok=True)
        stored_path = user_dir / safe_name
        stored_path.write_bytes(content)
        conn = _documents_db()
        conn.execute(
            'INSERT OR REPLACE INTO kb_documents '
            '(source_filename, stored_path, mimetype, ingested_at, chunks) VALUES (?, ?, ?, ?, ?)',
            (source_key, str(stored_path), mimetype, timestamp, len(chunks)),
        )
        conn.commit()
        conn.close()

        return len(chunks), source_key

    def ingest_document_by_page(self, content: bytes, filename: str, username: str) -> tuple:
        """Same Docling parse ingest_document() uses, but keeps page
        boundaries instead of collapsing the whole document into one
        markdown string first - export_to_markdown(page_no=N) pulls just
        that page's own text, one point per page in KB_PAGES_COLLECTION_NAME
        (not KB_COLLECTION_NAME - see that constant's own docstring for why
        pages and searchable chunks don't share a collection).

        For Scheduler's "send tonight's page" use case, not the knowledge-
        base search path - a document ingested here is NOT also searchable
        via retrieve_knowledge_base()/search_within_document(); call
        ingest_document() too if both are wanted. Doesn't shell out to
        _extract_pptx_markdown() - a slide deck has no meaningful "page N of
        a book" reading order, this is for genuinely paginated documents
        (PDFs, mainly).

        Returns (num_pages, source_filename), same shape as
        ingest_document()'s own return, for the same reason (a caller often
        wants to act on the just-ingested document right away)."""
        from docling.document_converter import DocumentConverter
        from docling_core.types.io import DocumentStream

        safe_user = _safe_filename(username)
        safe_name = _safe_filename(filename)
        source_key = f'{safe_user}/{safe_name}'

        converter = DocumentConverter()
        result = converter.convert(DocumentStream(name=safe_name, stream=io.BytesIO(content)))
        doc = result.document
        num_pages = doc.num_pages()
        if num_pages == 0:
            raise ValueError(f"Docling found no pages in '{filename}'")

        # Same delete-before-insert pattern ingest_document() uses - a
        # re-ingest must not leave stale pages from a previous version
        # sitting alongside the new ones. Guarded on collection_exists()
        # unlike ingest_document()'s own delete call - KB_COLLECTION_NAME
        # already has data from earlier sessions so that call never hit
        # this, but KB_PAGES_COLLECTION_NAME doesn't exist until the first
        # real insert below creates it, and delete() against a genuinely
        # missing collection 404s instead of being a harmless no-op.
        if self._qdrant_client.collection_exists(KB_PAGES_COLLECTION_NAME):
            self._qdrant_client.delete(
                collection_name=KB_PAGES_COLLECTION_NAME,
                points_selector=Filter(must=[FieldCondition(key='source_filename', match=MatchValue(value=source_key))]),
            )

        for page_no in range(1, num_pages + 1):
            page_text = doc.export_to_markdown(page_no=page_no)
            self.pages_index.insert(Document(
                text=page_text or ' ',  # a blank page still needs a non-empty embed input
                metadata={'source_filename': source_key, 'page_number': page_no},
            ))

        return num_pages, source_key

    def get_page(self, source_filename: str, page_number: int) -> Optional[str]:
        """One exact page's text, or None if that page (or the document
        itself) isn't on file - a metadata-filtered scroll(), not a
        similarity search, same reasoning as get_document_text()'s own
        docstring on why an exact lookup goes straight through the
        underlying Qdrant client instead of kb_index.as_retriever() (which
        only ever answers "most similar to a query," never "give me
        exactly this one thing")."""
        points, _ = self._qdrant_client.scroll(
            collection_name=KB_PAGES_COLLECTION_NAME,
            scroll_filter=Filter(must=[
                FieldCondition(key='source_filename', match=MatchValue(value=source_filename)),
                FieldCondition(key='page_number', match=MatchValue(value=page_number)),
            ]),
            limit=1,
            with_payload=True,
        )
        if not points:
            return None
        node_content = points[0].payload.get('_node_content')
        return json.loads(node_content).get('text', '') if node_content else None

    def list_page_ingested_books(self) -> List[str]:
        """Every distinct source_filename in KB_PAGES_COLLECTION_NAME - the
        page-ingested books AdiyanReader can actually read (ingest_book(),
        not ingest_document()). Deliberately separate from list_documents()
        (which only ever queries kb_documents, the chunk-ingested knowledge
        base): confirmed live this session that a book ingested via
        ingest_book() never gets a kb_documents row at all, so
        find_source_document() has zero visibility into it - resolving a
        "read me this book" reference has to search THIS collection
        instead, not the regular knowledge-base document list.

        Scrolls the whole collection client-side and dedupes source_filename
        values - no dedicated per-document index exists here (each point is
        one page, not one document), and the collection is small enough
        (one deployment's own uploaded books) that this is cheap."""
        if not self._qdrant_client.collection_exists(KB_PAGES_COLLECTION_NAME):
            return []
        seen: set = set()
        offset = None
        while True:
            points, offset = self._qdrant_client.scroll(
                collection_name=KB_PAGES_COLLECTION_NAME,
                limit=200,
                offset=offset,
                with_payload=['source_filename'],
            )
            for point in points:
                source_filename = point.payload.get('source_filename')
                if source_filename:
                    seen.add(source_filename)
            if offset is None:
                break
        return sorted(seen)

    def find_book_by_reference(self, query: str) -> Optional[str]:
        """Fuzzy-matches a free-text book reference (a title, a partial
        title) against the real page-ingested books on file, returning the
        exact source_filename or None if nothing matches well enough - the
        page-ingested equivalent of find_source_document(), which can't see
        these books at all (see list_page_ingested_books()'s own docstring).

        Plain string matching (difflib), not a semantic/LLM lookup - a book
        title is short and the candidate list is small, so a fuzzy string
        match against each book's own display name (the safe_filename with
        the <username>/ prefix and extension stripped, underscores back to
        spaces) is enough, and keeps this a real, inspectable lookup against
        actual data rather than another place an LLM could invent a key."""
        import difflib

        candidates = self.list_page_ingested_books()
        if not candidates:
            return None

        query_normalized = query.strip().lower()
        if not query_normalized:
            # Confirmed live: an empty/blank query is a substring of every
            # display name, so the exact/substring branch below would
            # "match" whatever book happens to come first in the whole
            # shared library and confidently start reading it - never a
            # deliberate choice. No reference at all means no match, full
            # stop, not "guess one."
            return None

        def _display_name(source_filename: str) -> str:
            basename = source_filename.split('/', 1)[-1]
            basename = re.sub(r'\.[A-Za-z0-9]+$', '', basename)
            return basename.replace('_', ' ').replace('-', ' ').strip().lower()

        display_to_source = {_display_name(c): c for c in candidates}

        # Exact/substring match first - a real book title mentioned in full
        # or as a clear substring shouldn't be left to difflib's fuzzier
        # scoring, which can be swayed by an unrelated but similarly-shaped
        # title.
        for display_name, source_filename in display_to_source.items():
            if query_normalized in display_name or display_name in query_normalized:
                return source_filename

        matches = difflib.get_close_matches(query_normalized, display_to_source.keys(), n=1, cutoff=0.4)
        return display_to_source[matches[0]] if matches else None

    def retrieve_knowledge_base(self, query: str, top_k: int = KB_DEFAULT_TOP_K) -> List[str]:
        """Return up to top_k relevant knowledge-base chunks, most relevant first.
        Global - not scoped to any one contact, unlike retrieve()."""
        retriever = self.kb_index.as_retriever(similarity_top_k=top_k)
        nodes = retriever.retrieve(query)
        return [n.node.get_content() for n in nodes]

    def search_within_document(
        self, source_filename: str, query: str, top_k: int = DOC_SEARCH_DEFAULT_TOP_K,
    ) -> List[Dict[str, Any]]:
        """Semantic search scoped to ONE already-known document, not the whole
        knowledge base - the fix for a confirmed-live gap: Analysis Agent's only
        way to get a document's content used to be get_document_text(), which
        concatenates EVERY chunk of a document into one giant string with no
        ranking at all. For a 167-chunk book that produced 400K+ characters,
        which the caller then truncates from the front before the ReAct loop
        ever reasons about it - so a passage deep in the middle was never seen,
        and the model fabricated a plausible-sounding answer instead.

        Reuses retrieve_knowledge_base()'s own retriever, just narrowed to one
        document's chunks via a metadata filter on source_filename. Returns
        each matching chunk's own text plus its score and chunk_index (not just
        the matched text alone) so a caller can cite exactly which part of the
        document an answer came from - mirrors what find_source_document()
        already keeps from a retrieved node rather than discarding it."""
        filters = MetadataFilters(filters=[
            MetadataFilter(key='source_filename', value=source_filename, operator=FilterOperator.EQ),
        ])
        retriever = self.kb_index.as_retriever(similarity_top_k=top_k, filters=filters)
        nodes = retriever.retrieve(query)
        return [
            {'text': n.node.get_content(), 'score': n.score, 'chunk_index': n.node.metadata.get('chunk_index')}
            for n in nodes
        ]

    def list_documents(self) -> List[str]:
        """Every source_filename currently on file in the raw-document index
        (kb_documents table), most recently ingested first - lets a caller
        see what's available without guessing or relying on a query
        happening to match, unlike find_source_document()'s similarity
        search below."""
        conn = _documents_db()
        rows = conn.execute('SELECT source_filename FROM kb_documents ORDER BY ingested_at DESC').fetchall()
        conn.close()
        return [row[0] for row in rows]

    def find_source_document(self, query: str) -> Optional[str]:
        """The source_filename of the single best-matching knowledge-base
        chunk for query, or None if the knowledge base has nothing at all,
        OR if the best match's own similarity score falls below
        SOURCE_MATCH_MIN_SCORE - a top-1 retriever always returns *something*
        as long as the knowledge base isn't empty, even when nothing in it is
        actually relevant (confirmed live: a Vizag-trip question against a
        knowledge base containing only an unrelated Atomic Habits PDF still
        got "matched" to it, and Analysis Agent's ReAct loop went on to
        answer from that instead of recognizing nothing relevant existed).
        Reuses the same chunk-level semantic search retrieve_knowledge_base()
        does - "which document answers this" is a different question from
        "what does it say," but the same search answers both, just keeping
        the winning chunk's metadata this time instead of discarding it."""
        retriever = self.kb_index.as_retriever(similarity_top_k=1)
        nodes = retriever.retrieve(query)
        if not nodes:
            return None
        best = nodes[0]
        if best.score is not None and best.score < SOURCE_MATCH_MIN_SCORE:
            return None
        return best.node.metadata.get('source_filename')

    def get_document_text(self, source_filename: str) -> Optional[str]:
        """Full text of a document ingested via ingest_document(), reconstructed by
        concatenating every stored chunk for source_filename in chunk_index order -
        not a similarity search (kb_index.as_retriever() only supports "most similar
        to a query," never "every chunk belonging to this exact document"), so this
        goes straight through the underlying Qdrant client instead. None if no chunks
        are on file for source_filename at all.

        LlamaIndex stores each chunk's own text inside its payload's _node_content
        field (a JSON-serialized TextNode), not as a plain top-level field - confirmed
        by inspecting a real stored point directly, not assumed."""
        points, _ = self._qdrant_client.scroll(
            collection_name=KB_COLLECTION_NAME,
            scroll_filter=Filter(
                must=[FieldCondition(key='source_filename', match=MatchValue(value=source_filename))]
            ),
            limit=10000,
            with_payload=True,
        )
        if not points:
            return None
        ordered = sorted(points, key=lambda p: p.payload.get('chunk_index', 0))
        texts = []
        for point in ordered:
            node_content = point.payload.get('_node_content')
            if node_content:
                texts.append(json.loads(node_content).get('text', ''))
        return '\n\n'.join(texts)

    def get_document(self, source_filename: str) -> Optional[Dict[str, str]]:
        """Raw bytes (base64) plus mimetype for a document previously ingested via
        ingest_document(), keyed by the same source_filename its chunks carry in Qdrant -
        or None if no such document is on file (never actually ingested, or ingested
        before this raw-storage index existed, or its stored file went missing).

        The returned 'filename' is just the original basename, not the internal
        <username>/<filename> key - that composite form is this index's own identity
        for avoiding cross-user collisions, not something that should show up as a
        WhatsApp attachment's displayed name."""
        conn = _documents_db()
        row = conn.execute(
            'SELECT stored_path, mimetype FROM kb_documents WHERE source_filename = ?',
            (source_filename,),
        ).fetchone()
        conn.close()
        if row is None:
            return None
        stored_path, mimetype = row
        path = Path(stored_path)
        if not path.exists():
            return None
        return {
            'filename': source_filename.rsplit('/', 1)[-1],
            'mimetype': mimetype or 'application/octet-stream',
            'content_b64': base64.b64encode(path.read_bytes()).decode('ascii'),
        }


def get_memory_index(qdrant_url: str, ollama_url: str) -> Optional[MemoryIndex]:
    """Cached factory - returns None (not an error) if Qdrant/Ollama aren't reachable,
    so callers can treat memory as an optional enhancement rather than a hard dependency."""
    key = (qdrant_url, ollama_url)
    if key not in _instances:
        try:
            _instances[key] = MemoryIndex(qdrant_url, ollama_url)
        except Exception as e:
            logger.warning(f"⚠️  Memory index unavailable ({e}) - continuing without it")
            _instances[key] = None
    return _instances[key]
